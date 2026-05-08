from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──────────────────────────────────────────
BRAND       = RGBColor(0x1A, 0x73, 0xE8)
BRAND_LIGHT = RGBColor(0xE8, 0xF0, 0xFE)
DARK        = RGBColor(0x1C, 0x1E, 0x21)
SUB         = RGBColor(0x65, 0x67, 0x6B)
GREEN       = RGBColor(0x1B, 0x9E, 0x4B)
RED         = RGBColor(0xE5, 0x39, 0x35)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG          = RGBColor(0xF0, 0xF2, 0xF5)
GOLD        = RGBColor(0xFF, 0xA0, 0x00)
TEAL        = RGBColor(0x00, 0x96, 0x88)

# ── Helpers ─────────────────────────────────────────
def add_bg(slide, color=BRAND):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def set_text(shape, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Microsoft JhengHei'):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Microsoft JhengHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=DARK, spacing=Pt(8), font_name='Microsoft JhengHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = clr if clr else color
        p.font.name = font_name
        p.space_after = spacing
    return txBox

def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=SUB, alignment=PP_ALIGN.RIGHT)

def add_accent_bar(slide, left, top, height, color=BRAND):
    add_rect(slide, left, top, Inches(0.06), height, color)

TOTAL_SLIDES = 14

# ════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_rect(slide, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), BRAND)

# Left accent
add_rect(slide, Inches(0), Inches(0), Inches(0.35), Inches(7.5), BRAND)

# Main title area
add_text_box(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(1.2),
             "StockBoard 股票交易論壇", font_size=40, bold=True, color=BRAND)
add_text_box(slide, Inches(1.2), Inches(2.7), Inches(10), Inches(0.8),
             "全端開發實戰報告", font_size=28, bold=False, color=DARK)

# Decorative line
add_rect(slide, Inches(1.2), Inches(3.7), Inches(3), Inches(0.04), BRAND)

# Subtitle info
info_lines = [
    ("技術堆疊：FastAPI + SQLAlchemy + Vanilla JS + Chart.js", False, SUB),
    ("部署平台：Render.com (PostgreSQL)", False, SUB),
    ("專案性質：股票討論、查詢、交易紀錄一平台", False, SUB),
]
add_multi_text(slide, Inches(1.2), Inches(4.0), Inches(8), Inches(2), info_lines, font_size=16)

# Tech badges
techs = ["FastAPI", "SQLAlchemy", "yfinance", "Chart.js", "SQLite / PostgreSQL", "bcrypt"]
for i, tech in enumerate(techs):
    left = Inches(1.2 + i * 2.0)
    badge = add_rounded_rect(slide, left, Inches(5.8), Inches(1.7), Inches(0.5), BRAND_LIGHT)
    set_text(badge, tech, font_size=12, bold=True, color=BRAND, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.6), Inches(5), Inches(0.8),
             "目 錄", font_size=32, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.35), Inches(2), Inches(0.04), BRAND)

agenda_items = [
    ("01", "專案概述與動機"),
    ("02", "系統架構設計"),
    ("03", "資料庫設計 (ERD)"),
    ("04", "後端 API 設計"),
    ("05", "前端 UI/UX 設計"),
    ("06", "核心功能展示"),
    ("07", "技術亮點與創新"),
    ("08", "部署與維運"),
    ("09", "安全性分析"),
    ("10", "遇到的挑戰與解法"),
    ("11", "未來改進方向"),
    ("12", "總結與心得"),
]
for i, (num, title) in enumerate(agenda_items):
    row = i % 6
    col = i // 6
    left = Inches(1.0 + col * 5.8)
    top = Inches(1.8 + row * 0.85)
    add_text_box(slide, left, top, Inches(0.6), Inches(0.5),
                 num, font_size=22, bold=True, color=BRAND)
    add_text_box(slide, left + Inches(0.7), top, Inches(4.5), Inches(0.5),
                 title, font_size=16, bold=False, color=DARK)

add_slide_number(slide, 2, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 3: Project Overview
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(6), Inches(0.7),
             "01  專案概述", font_size=28, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.15), Inches(2.5), Inches(0.04), BRAND)

# Left card
card1 = add_rounded_rect(slide, Inches(1.0), Inches(1.6), Inches(5.5), Inches(2.5), BRAND_LIGHT)
card1.text_frame.margin_left = Inches(0.3)
card1.text_frame.margin_top = Inches(0.2)
tf = card1.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "專案目標"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BRAND
tf.paragraphs[0].font.name = 'Microsoft JhengHei'
for t in ["打造一個整合股票討論、即時查詢、交易紀錄的全端平台",
          "讓投資者可以交流投資心得、追蹤股價走勢、記錄交易歷史",
          "實現前後端分離架構，具備良好的擴展性與部署能力"]:
    p = tf.add_paragraph()
    p.text = f"▸ {t}"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft JhengHei'
    p.space_before = Pt(8)

# Right card
card2 = add_rounded_rect(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.5), RGBColor(0xE8, 0xF5, 0xE9))
card2.text_frame.margin_left = Inches(0.3)
card2.text_frame.margin_top = Inches(0.2)
tf = card2.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "核心特色"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = GREEN
tf.paragraphs[0].font.name = 'Microsoft JhengHei'
for t in ["社群討論板：分類看板、發文、留言互動",
          "股票即時查詢：串接 yfinance API，圖表化呈現",
          "交易紀錄管理：買入/賣出記錄、歷史查詢",
          "使用者個人檔案：文章統計、追蹤功能"]:
    p = tf.add_paragraph()
    p.text = f"▸ {t}"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft JhengHei'
    p.space_before = Pt(8)

# Bottom: Tech stack overview
add_text_box(slide, Inches(1.0), Inches(4.5), Inches(5), Inches(0.5),
             "技術堆疊總覽", font_size=20, bold=True, color=DARK)

stack_data = [
    ("後端", "FastAPI (Python)", BRAND),
    ("資料庫", "SQLite (dev) / PostgreSQL (prod)", GREEN),
    ("前端", "Vanilla HTML/CSS/JS + Chart.js", TEAL),
    ("部署", "Render.com (免費層級)", GOLD),
]
for i, (label, desc, clr) in enumerate(stack_data):
    top = Inches(5.1 + i * 0.55)
    add_rounded_rect(slide, Inches(1.0), top, Inches(1.2), Inches(0.42), clr)
    lb = add_text_box(slide, Inches(1.0), top, Inches(1.2), Inches(0.42),
                      label, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.4), top, Inches(10), Inches(0.42),
                 desc, font_size=14, color=DARK)

add_slide_number(slide, 3, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 4: System Architecture
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(6), Inches(0.7),
             "02  系統架構設計", font_size=28, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.15), Inches(2.5), Inches(0.04), BRAND)

# Architecture diagram boxes
layers = [
    ("前端層 (Frontend)", "HTML5 + CSS3 + Vanilla JavaScript\nChart.js (圖表渲染)\nLocalStorage (Session 管理)", BRAND),
    ("API 層 (Backend)", "FastAPI (ASGI)\nRESTful API 設計\nPydantic 資料驗證\nbcrypt 密碼雜湊", GREEN),
    ("資料層 (Database)", "SQLAlchemy ORM\nSQLite (開發環境)\nPostgreSQL (生產環境)\nyfinance (外部 API)", TEAL),
]
for i, (title, desc, clr) in enumerate(layers):
    top = Inches(1.7 + i * 1.75)
    # Left box
    box = add_rounded_rect(slide, Inches(1.0), top, Inches(2.8), Inches(1.4), clr)
    set_text(box, "", font_size=1)
    box.text_frame.margin_left = Inches(0.2)
    box.text_frame.margin_top = Inches(0.25)
    box.text_frame.word_wrap = True
    tf = box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = 'Microsoft JhengHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Arrow
    if i < 2:
        add_text_box(slide, Inches(2.1), top + Inches(1.45), Inches(0.5), Inches(0.3),
                     "▼", font_size=18, bold=True, color=clr, alignment=PP_ALIGN.CENTER)

    # Right box
    box2 = add_rounded_rect(slide, Inches(4.5), top, Inches(7.5), Inches(1.4), BG)
    set_text(box2, "", font_size=1)
    box2.text_frame.margin_left = Inches(0.3)
    box2.text_frame.margin_top = Inches(0.2)
    box2.text_frame.word_wrap = True
    tf = box2.text_frame
    tf.paragraphs[0].text = desc
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = DARK
    tf.paragraphs[0].font.name = 'Microsoft JhengHei'

# File structure
add_text_box(slide, Inches(1.0), Inches(6.8), Inches(11), Inches(0.5),
             "專案結構：main.py (API) │ models.py (ORM) │ database.py (DB Config) │ index.html (SPA)",
             font_size=12, color=SUB, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 4, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 5: Database Design
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(6), Inches(0.7),
             "03  資料庫設計", font_size=28, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.15), Inches(2.5), Inches(0.04), BRAND)

# Tables
tables = [
    ("Users (使用者)", [("id", "Integer", "PK"), ("username", "String", "UNIQUE, INDEX"), ("password", "String(500)", "bcrypt hashed")], BRAND),
    ("TradeRecords (交易)", [("id", "Integer", "PK"), ("user_id", "Integer", "FK → users.id"), ("stock_symbol", "String", ""), ("trade_type", "Enum", "BUY / SELL"), ("price", "Float", ""), ("quantity", "Integer", ""), ("trade_date", "DateTime", "UTC")], GREEN),
    ("Posts (文章)", [("id", "Integer", "PK"), ("user_id", "Integer", "FK → users.id"), ("title", "String", ""), ("content", "String", ""), ("created_at", "DateTime", "UTC")], TEAL),
    ("Comments (留言)", [("id", "Integer", "PK"), ("post_id", "Integer", "FK → posts.id"), ("user_id", "Integer", "FK → users.id"), ("content", "String", ""), ("created_at", "DateTime", "UTC")], GOLD),
]
for i, (title, cols, clr) in enumerate(tables):
    col_idx = i % 4
    row_idx = i // 4
    left = Inches(0.8 + col_idx * 3.1)
    top = Inches(1.6 + row_idx * 3.0)

    # Header
    hdr = add_rounded_rect(slide, left, top, Inches(2.85), Inches(0.45), clr)
    set_text(hdr, "", font_size=1)
    hdr.text_frame.margin_left = Inches(0.15)
    tf = hdr.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = 'Microsoft JhengHei'

    # Columns
    for j, (name, typ, extra) in enumerate(cols):
        y = top + Inches(0.48 + j * 0.33)
        add_text_box(slide, left + Inches(0.1), y, Inches(1.0), Inches(0.3),
                     name, font_size=10, bold=True, color=DARK)
        add_text_box(slide, left + Inches(1.1), y, Inches(0.8), Inches(0.3),
                     typ, font_size=10, color=SUB)
        if extra:
            add_text_box(slide, left + Inches(1.9), y, Inches(0.85), Inches(0.3),
                         extra, font_size=9, color=BRAND)

add_slide_number(slide, 5, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 6: API Design
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(6), Inches(0.7),
             "04  後端 API 設計", font_size=28, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.15), Inches(2.5), Inches(0.04), BRAND)

# API table headers
headers = ["端點 (Endpoint)", "方法", "功能說明"]
col_widths = [Inches(3.5), Inches(1.0), Inches(6.5)]
x_start = Inches(1.0)
y = Inches(1.5)

for j, (hdr, w) in enumerate(zip(headers, col_widths)):
    x = x_start + sum(cw.inches for cw in col_widths[:j]) * Inches(1)
    box = add_rounded_rect(slide, x, y, w, Inches(0.45), BRAND)
    set_text(box, hdr, font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# API endpoints
endpoints = [
    ("/register/", "POST", "使用者註冊：檢查重複、bcrypt 雜湊密碼、建立帳號", RGBColor(0xE3, 0xF2, 0xFD)),
    ("/login/", "POST", "使用者登入：驗證密碼、回傳 user_id 與 username", RGBColor(0xE3, 0xF2, 0xFD)),
    ("/trades/", "POST", "建立交易紀錄：驗證使用者、Enum 轉換、寫入資料庫", RGBColor(0xE8, 0xF5, 0xE9)),
    ("/trades/{user_id}", "GET", "查詢使用者所有交易紀錄", RGBColor(0xE8, 0xF5, 0xE9)),
    ("/posts/", "POST", "建立新文章：驗證使用者、寫入文章資料", RGBColor(0xE0, 0xF2, 0xF1)),
    ("/posts/", "GET", "取得所有文章（含巢狀留言），依時間倒序", RGBColor(0xE0, 0xF2, 0xF1)),
    ("/comments/", "POST", "對文章新增留言：驗證使用者與文章存在性", RGBColor(0xE0, 0xF2, 0xF1)),
    ("/stock/{symbol}", "GET", "即時股價：透過 yfinance 取得當前價格、高低價", RGBColor(0xFF, 0xE0, 0xB2)),
    ("/stock/{symbol}/history", "GET", "歷史走勢：支援 1d/5d/1mo/3mo/6mo/1y 區間", RGBColor(0xFF, 0xE0, 0xB2)),
]
for i, (ep, method, desc, bg_clr) in enumerate(endpoints):
    y = Inches(2.0 + i * 0.58)
    # Alternate background
    if i % 2 == 1:
        add_rect(slide, Inches(0.9), y, Inches(11.5), Inches(0.55), BG)

    add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(3.3), Inches(0.45),
                 ep, font_size=12, bold=True, color=BRAND)

    method_clr = GREEN if "GET" in method else RGBColor(0xFF, 0x6D, 0x00)
    m_box = add_rounded_rect(slide, Inches(4.6), y + Inches(0.05), Inches(0.8), Inches(0.35), method_clr)
    set_text(m_box, method, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(5.6), y + Inches(0.05), Inches(6.5), Inches(0.45),
                 desc, font_size=12, color=DARK)

add_slide_number(slide, 6, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 7: Frontend UI/UX
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(6), Inches(0.7),
             "05  前端 UI/UX 設計", font_size=28, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.15), Inches(2.5), Inches(0.04), BRAND)

# Left: Layout overview
add_text_box(slide, Inches(1.0), Inches(1.5), Inches(5), Inches(0.5),
             "SPA 架構設計", font_size=20, bold=True, color=DARK)

layout_items = [
    ("Topbar", "固定頂部導航列：Logo、搜尋列、使用者資訊"),
    ("Sidebar", "固定左側選單：功能導航、看板分類、快捷發文"),
    ("Main Content", "可切換頁面區域：討論板 / 股票查詢 / 交易 / 個人檔案"),
    ("Modal", "彈出式對話框：發文表單、編輯個人資料"),
    ("Auth Overlay", "登入/註冊浮層：Backdrop blur 動畫效果"),
]
for i, (title, desc) in enumerate(layout_items):
    y = Inches(2.1 + i * 0.7)
    badge = add_rounded_rect(slide, Inches(1.0), y, Inches(1.5), Inches(0.4), BRAND_LIGHT)
    set_text(badge, title, font_size=11, bold=True, color=BRAND, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.7), y, Inches(7), Inches(0.4),
                 desc, font_size=13, color=DARK)

# Right: Design patterns
card = add_rounded_rect(slide, Inches(7.5), Inches(1.5), Inches(5.0), Inches(4.8), BG)
card.text_frame.margin_left = Inches(0.3)
card.text_frame.margin_top = Inches(0.2)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "UI 設計模式"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BRAND
tf.paragraphs[0].font.name = 'Microsoft JhengHei'

patterns = [
    "CSS 變數主題系統 (--brand, --bg, --card)",
    "CSS 動畫：authIn、spin、toast slide-in",
    "Template Literals 動態渲染文章/留言",
    "esc() XSS 防護函數",
    "Toast 通知系統 (success/error/info)",
    "伺服器喚醒指示器 (Render 冷啟動)",
    "RWD 響應式設計 (mobile: sidebar 隱藏)",
    "Noto Sans TC 中文字型",
]
for p in patterns:
    p_obj = tf.add_paragraph()
    p_obj.text = f"▸ {p}"
    p_obj.font.size = Pt(13)
    p_obj.font.color.rgb = DARK
    p_obj.font.name = 'Microsoft JhengHei'
    p_obj.space_before = Pt(8)

# Bottom: JS architecture
add_text_box(slide, Inches(1.0), Inches(6.6), Inches(11), Inches(0.5),
             "前端模組化：Auth │ Posts │ Stock │ Trades │ Profile │ Board Filter 六大功能模組",
             font_size=13, color=SUB, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 7, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 8: Core Features
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(6), Inches(0.7),
             "06  核心功能展示", font_size=28, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.15), Inches(2.5), Inches(0.04), BRAND)

features = [
    ("討論板", "發文、留言互動、看板分類篩選（台股/美股/加密貨幣/新手）", BRAND, "0xE8F0FE"),
    ("股票查詢", "即時股價、K線圖表、多時間區間切換（5天~1年）", GREEN, "0xE8F5E9"),
    ("交易紀錄", "買入/賣出記錄、歷史查詢、圖標化顯示", TEAL, "0xE0F2F1"),
    ("個人檔案", "文章統計、追蹤功能、個人資料編輯、查看他人檔案", GOLD, "0xFFF3E0"),
]
for i, (title, desc, clr, bg_hex) in enumerate(features):
    col = i % 2
    row = i // 2
    left = Inches(1.0 + col * 5.8)
    top = Inches(1.6 + row * 2.8)

    card = add_rounded_rect(slide, left, top, Inches(5.4), Inches(2.4), BG)

    # Title bar
    title_bar = add_rounded_rect(slide, left + Inches(0.05), top + Inches(0.05), Inches(5.3), Inches(0.55), clr)
    set_text(title_bar, "", font_size=1)
    title_bar.text_frame.margin_left = Inches(0.25)
    title_bar.text_frame.margin_top = Inches(0.1)
    tf = title_bar.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(17)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = 'Microsoft JhengHei'

    # Description
    add_multi_text(slide, left + Inches(0.3), top + Inches(0.8), Inches(4.8), Inches(1.4),
                   [(desc, False, DARK)], font_size=14)

add_slide_number(slide, 8, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 9: Technical Highlights
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.7),
             "07  技術亮點", font_size=28, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.15), Inches(2.5), Inches(0.04), BRAND)

highlights = [
    ("字元級 bcrypt 截斷",
     "bcrypt 限制 72 bytes，本專案實作了字元層級截斷邏輯，避免中文字被切斷。\n一般做法直接截斷 byte 72，會破壞多字元 UTF-8 字元。",
     BRAND),
    ("雙資料庫引擎配置",
     "database.py 自動偵測 SQLite vs PostgreSQL，配置正確的 engine 參數。\n處理 postgres:// 到 postgresql:// 的 URL scheme 轉換（Render 部署需要）。",
     GREEN),
    ("取代棄用的 datetime.utcnow()",
     "models.py 使用 datetime.now(timezone.utc).replace(tzinfo=None) 取代已棄用的 utcnow()，\n同時維持 SQLite 的 naive datetime 相容性。",
     TEAL),
    ("Render 冷啟動處理",
     "fetchWT() 包裝 fetch 帶 90 秒逾時，4 秒後顯示「伺服器喚醒中」spinner。\n解決 Render 免費層級休眠後首次請求延遲的問題。",
     GOLD),
    ("台股代號正規化",
     "normTicker() 自動轉換台股代號：2330 → 2330.TW，TW2330 → 2330.TW。\n支援 Yahoo Finance 的 TW 前綴格式與標準 4-5 位數字代號。",
     RED),
]
for i, (title, desc, clr) in enumerate(highlights):
    top = Inches(1.6 + i * 1.1)

    num_badge = add_rounded_rect(slide, Inches(1.0), top, Inches(0.5), Inches(0.5), clr)
    set_text(num_badge, str(i+1), font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.7), top, Inches(3.5), Inches(0.4),
                 title, font_size=15, bold=True, color=DARK)
    add_text_box(slide, Inches(1.7), top + Inches(0.4), Inches(10.5), Inches(0.6),
                 desc, font_size=11, color=SUB)

add_slide_number(slide, 9, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 10: Deployment & Operations
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.7),
             "08  部署與維運", font_size=28, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.15), Inches(2.5), Inches(0.04), BRAND)

# Deployment flow
flow_items = [
    ("本地開發", "SQLite + uvicorn\nlocalhost:8000", BRAND),
    ("版本控制", "Git 提交\n. gitignore 排除\n.pyc 和 .db", GREEN),
    ("雲端部署", "Render.com\n自動偵測 requirements.txt\nPostgreSQL 注入", TEAL),
    ("前端部署", "index.html\n獨立靜態檔案\nCDN 載入 Chart.js", GOLD),
]
for i, (title, desc, clr) in enumerate(flow_items):
    left = Inches(1.0 + i * 3.0)
    top = Inches(1.6)

    box = add_rounded_rect(slide, left, top, Inches(2.6), Inches(2.2), BG)
    box.text_frame.margin_left = Inches(0.2)
    box.text_frame.margin_top = Inches(0.2)
    tf = box.text_frame
    tf.word_wrap = True

    hdr = add_rounded_rect(slide, left + Inches(0.05), top + Inches(0.05), Inches(2.5), Inches(0.5), clr)
    set_text(hdr, title, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(0.2), top + Inches(0.7), Inches(2.2), Inches(1.3),
                 desc, font_size=12, color=DARK)

    if i < 3:
        add_text_box(slide, left + Inches(2.6), top + Inches(0.9), Inches(0.4), Inches(0.4),
                     "→", font_size=20, bold=True, color=clr, alignment=PP_ALIGN.CENTER)

# Key points
card = add_rounded_rect(slide, Inches(1.0), Inches(4.2), Inches(11.3), Inches(2.8), BRAND_LIGHT)
card.text_frame.margin_left = Inches(0.3)
card.text_frame.margin_top = Inches(0.2)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "部署關鍵設定"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BRAND
tf.paragraphs[0].font.name = 'Microsoft JhengHei'

deploy_points = [
    "DATABASE_URL 環境變數：本地預設 SQLite，Render 自動注入 PostgreSQL 連線字串",
    "Lifespan 管理：FastAPI 啟動時自動呼叫 create_tables() 建立資料表",
    "CORS 設定：allow_origins=[\"*\"] 允許所有來源（開發環境用）",
    "requirements.txt 包含：fastapi, uvicorn, sqlalchemy, psycopg2-binary, bcrypt, yfinance, pydantic",
    "啟動命令：uvicorn main:app --host 0.0.0.0 --port 8000",
    "生產環境 URL: https://web-design-67t2.onrender.com",
]
for p in deploy_points:
    p_obj = tf.add_paragraph()
    p_obj.text = f"▸ {p}"
    p_obj.font.size = Pt(13)
    p_obj.font.color.rgb = DARK
    p_obj.font.name = 'Microsoft JhengHei'
    p_obj.space_before = Pt(4)

add_slide_number(slide, 10, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 11: Security Analysis
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.7),
             "09  安全性分析", font_size=28, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.15), Inches(2.5), Inches(0.04), BRAND)

# Good practices
add_text_box(slide, Inches(1.0), Inches(1.5), Inches(5), Inches(0.5),
             "已實作的安全措施", font_size=18, bold=True, color=GREEN)

good_items = [
    "bcrypt 密碼雜湊（含 salt），不使用明文儲存",
    "字元層級 bcrypt 截斷，避免中文字元損壞",
    "esc() XSS 防護函數，過濾 &, <, >, \", '",
    "SQLAlchemy ORM 防護 SQL Injection",
    "HTTPException 統一錯誤處理 + db.rollback()",
    "Pydantic 輸入資料驗證（型別檢查）",
]
for i, item in enumerate(good_items):
    y = Inches(2.1 + i * 0.5)
    add_text_box(slide, Inches(1.5), y, Inches(0.3), Inches(0.3),
                 "✓", font_size=14, bold=True, color=GREEN)
    add_text_box(slide, Inches(1.9), y, Inches(5), Inches(0.3),
                 item, font_size=13, color=DARK)

# Areas for improvement
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5), Inches(0.5),
             "可改進項目", font_size=18, bold=True, color=RED)

improve_items = [
    "缺少 JWT/Token 認證機制",
    "API 信任前端傳送的 user_id",
    "CORS allow_origins=[\"*\"] 過於寬鬆",
    "無速率限制 (Rate Limiting)",
    "個人資料 (bio/following) 存在 localStorage",
    "無 CSRF Token 防護",
]
for i, item in enumerate(improve_items):
    y = Inches(2.1 + i * 0.5)
    add_text_box(slide, Inches(7.5), y, Inches(0.3), Inches(0.3),
                 "△", font_size=14, bold=True, color=RED)
    add_text_box(slide, Inches(7.9), y, Inches(5), Inches(0.3),
                 item, font_size=13, color=DARK)

# Recommendation card
rec = add_rounded_rect(slide, Inches(1.0), Inches(5.3), Inches(11.3), Inches(1.6), RGBColor(0xFF, 0xF3, 0xE0))
rec.text_frame.margin_left = Inches(0.3)
rec.text_frame.margin_top = Inches(0.15)
tf = rec.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "建議改進方向"
tf.paragraphs[0].font.size = Pt(17)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0xE6, 0x51, 0x00)
tf.paragraphs[0].font.name = 'Microsoft JhengHei'

rec_items = [
    "導入 JWT 驗證：login 後回發 access token，後續 API 從 Header 驗證",
    "限制 CORS：設定 allow_origins 為特定網域",
    "加入 Rate Limiting：使用 slowapi 防止暴力攻擊",
    "個人資料遷移至後端：bio/following 存入資料庫",
]
for item in rec_items:
    p = tf.add_paragraph()
    p.text = f"▸ {item}"
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft JhengHei'
    p.space_before = Pt(3)

add_slide_number(slide, 11, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 12: Challenges & Solutions
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.7),
             "10  遇到的挑戰與解法", font_size=28, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.15), Inches(2.5), Inches(0.04), BRAND)

challenges = [
    ("Render 免費層冷啟動",
     "問題：Render 免費層休眠後首次請求需 30-50 秒，使用者看到逾時",
     "解法：fetchWT() 實作 90 秒逾時 + 4 秒後顯示「喚醒中」spinner，提供明確回饋",
     BRAND),
    ("bcrypt 中文字元處理",
     "問題：密碼超過 72 bytes 需截斷，但中文字元是 3 bytes，直接截斷會損壞",
     "解法：實作 _trim_password() 以字元為單位遍歷，確保在完整字元邊界處截斷",
     GREEN),
    ("資料庫跨平台相容",
     "問題：SQLite 需要 check_same_thread=False，PostgreSQL 不需要且 URL 格式不同",
     "解法：database.py 自動偵測 URL prefix，動態配置 engine 參數",
     TEAL),
    ("台股代號轉換",
     "問題：Yahoo Finance 台股需加 .TW 後綴，但使用者習慣輸入純數字",
     "解法：normTicker() 正規化函數，自動補齊/移除 TW 前綴",
     GOLD),
]
for i, (title, problem, solution, clr) in enumerate(challenges):
    col = i % 2
    row = i // 2
    left = Inches(1.0 + col * 5.8)
    top = Inches(1.6 + row * 2.8)

    card = add_rounded_rect(slide, left, top, Inches(5.4), Inches(2.5), BG)

    # Title
    hdr = add_rounded_rect(slide, left + Inches(0.05), top + Inches(0.05), Inches(5.3), Inches(0.45), clr)
    set_text(hdr, title, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(0.2), top + Inches(0.6), Inches(5.0), Inches(0.8),
                 problem, font_size=12, color=RED)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.4), Inches(5.0), Inches(0.9),
                 solution, font_size=12, color=GREEN)

add_slide_number(slide, 12, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 13: Future Improvements
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.7),
             "11  未來改進方向", font_size=28, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.15), Inches(2.5), Inches(0.04), BRAND)

improvements = [
    ("安全性強化", ["JWT Token 認證機制", "CORS 網域限制", "Rate Limiting", "CSRF Token"], BRAND),
    ("功能擴充", ["文章按讚/收藏", "即時通知系統", "投资组合分析儀表板", "搜尋功能強化"], GREEN),
    ("使用者體驗", ["深色模式 (Dark Mode)", "RWD 行動版優化", "圖片上傳支援", "Markdown 編輯器"], TEAL),
    ("效能優化", ["資料庫查詢快取 (Redis)", "分頁載入 (Pagination)", "WebSocket 即時更新", "CDN 靜態資源加速"], GOLD),
    ("資料層", ["個人檔案後端化", "文章分類標籤系統", "交易統計圖表", "匯出交易紀錄"], RED),
]
for i, (category, items, clr) in enumerate(improvements):
    col = i % 3
    row = i // 3
    left = Inches(1.0 + col * 3.9)
    top = Inches(1.6 + row * 3.0)

    hdr = add_rounded_rect(slide, left, top, Inches(3.6), Inches(0.5), clr)
    set_text(hdr, category, font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        y = top + Inches(0.65 + j * 0.42)
        add_text_box(slide, left + Inches(0.3), y, Inches(0.3), Inches(0.3),
                     "●", font_size=10, color=clr)
        add_text_box(slide, left + Inches(0.6), y, Inches(2.8), Inches(0.3),
                     item, font_size=13, color=DARK)

add_slide_number(slide, 13, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════
# SLIDE 14: Summary
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BRAND)
add_accent_bar(slide, Inches(0), Inches(0), Inches(7.5))

add_text_box(slide, Inches(1.0), Inches(0.8), Inches(8), Inches(0.8),
             "總結", font_size=36, bold=True, color=BRAND)
add_rect(slide, Inches(1.0), Inches(1.55), Inches(3), Inches(0.04), BRAND)

# Summary points
summary_items = [
    "完成了一個全端股票交易論壇，整合社群、查詢、交易三大核心功能",
    "實踐了 FastAPI + SQLAlchemy ORM + Vanilla JS 的現代化技術堆疊",
    "成功部署至 Render.com，實現本地開發與生產環境的無縫切換",
    "深入理解了 RESTful API 設計、資料庫關聯、前後端分離架構",
    "學會處理實際部署挑戰：冷啟動、跨資料庫相容、字元編碼",
]
for i, item in enumerate(summary_items):
    y = Inches(2.0 + i * 0.6)
    num = add_rounded_rect(slide, Inches(1.0), y, Inches(0.45), Inches(0.45), BRAND)
    set_text(num, str(i+1), font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.7), y + Inches(0.05), Inches(10), Inches(0.4),
                 item, font_size=15, color=DARK)

# Thank you card
ty_card = add_rounded_rect(slide, Inches(3.5), Inches(5.2), Inches(6.3), Inches(1.5), BRAND_LIGHT)
ty_card.text_frame.margin_left = Inches(0.3)
ty_card.text_frame.margin_top = Inches(0.15)
tf = ty_card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "感謝聆聽"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BRAND
tf.paragraphs[0].font.name = 'Microsoft JhengHei'
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "StockBoard 股票交易論壇 專案報告"
p.font.size = Pt(14)
p.font.color.rgb = SUB
p.font.name = 'Microsoft JhengHei'
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(8)

add_slide_number(slide, 14, TOTAL_SLIDES)

# ── Save ────────────────────────────────────────────────
output_path = os.path.join(os.getcwd(), "StockBoard_專案報告_v2.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
